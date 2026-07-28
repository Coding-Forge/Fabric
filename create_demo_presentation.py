"""
Create a professional PowerPoint presentation from Fabric Monitor screenshots.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Configuration
IMAGE_DIR = Path(__file__).parent / "images"
OUTPUT_PPTX = Path(__file__).parent / "Fabric_Monitor_Demo.pptx"

# Define slide content: (image_filename, title, description)
SLIDES = [
    (
        None,
        "Fabric Monitor",
        "Comprehensive Power BI & Fabric Estate Monitoring\n\nSupports Azure Commercial, GCC, GCC High"
    ),
    (
        "admin-portal-settings.png",
        "Setup: Admin Portal Configuration",
        "Grant Service Principal permissions in Fabric Admin Portal\nEnable API access for monitoring and data extraction"
    ),
    (
        "Service-Principal-API-Permissions.png",
        "Setup: API Permissions",
        "Configure Graph API and Power BI permissions\nApplication-level permissions for headless monitoring"
    ),
    (
        None,
        "Monitoring Capabilities",
        "Real-time visibility into your Power BI and Fabric estate"
    ),
    (
        "Workspace_Artifacts.png",
        "Workspace Inventory",
        "Complete catalog of all workspaces and artifacts\nTrack reports, datasets, dashboards, and dataflows"
    ),
    (
        "Catalog_Governance.png",
        "Governance & Classification",
        "Monitor data classification and governance policies\nEnsure compliance across your enterprise"
    ),
    (
        "Activity_Operations.png",
        "Activity & Operations",
        "Track user activity, refresh operations, and usage patterns\nIdentify trends and monitor platform health"
    ),
    (
        "PBI_Audience_Usage.png",
        "Audience & Usage Analytics",
        "Understand which audiences are consuming content\nMeasure engagement and adoption metrics"
    ),
    (
        "Users_Access.png",
        "User Access & Security",
        "Monitor user permissions and access levels\nEnsure proper authorization across workspaces"
    ),
    (
        "User_Inactivity.png",
        "User Inactivity Tracking",
        "Identify inactive users and unused capacity\nOptimize licensing and capacity allocation"
    ),
    (
        "RiskAnomalies.png",
        "Risk & Anomaly Detection",
        "Detect unusual patterns and potential security risks\nProactive monitoring for estate health"
    ),
    (
        "Semantic_Model_Governance.png",
        "Semantic Model Governance",
        "Monitor data models, dependencies, and quality\nTrack lineage and data governance metrics"
    ),
    (
        "Audit_Overview.png",
        "Audit Overview Dashboard",
        "Comprehensive view of all monitoring activities\nCentralized dashboard for estate oversight"
    ),
    (
        "Customer_Audit_Overview.png",
        "Customer-Facing Audit View",
        "Present insights to stakeholders and leadership\nTransparent reporting of platform metrics"
    ),
    (
        "Drillthrough_Details.png",
        "Detailed Drill-Through Analysis",
        "Dive deep into specific items for root cause analysis\nUnderstand the 'why' behind the metrics"
    ),
    (
        None,
        "Multi-Cloud Support",
        "Azure Commercial • GCC • GCC High\n\nSeamless monitoring across all government and commercial clouds"
    ),
    (
        None,
        "Key Benefits",
        "✓ Real-time visibility across entire estate\n✓ Governance and compliance monitoring\n✓ Usage and adoption analytics\n✓ Proactive risk detection"
    ),
    (
        None,
        "Get Started Today",
        "Deploy Fabric Monitor to gain complete visibility\ninto your Power BI and Fabric environment"
    ),
]

# Colors
DARK_BLUE = RGBColor(25, 35, 65)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_TEXT = RGBColor(25, 35, 65)
LIGHT_BLUE = RGBColor(176, 190, 197)


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(9), Inches(2)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, image_path, title, description):
    """Add a slide with image and description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    
    # Add image
    try:
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5), Inches(1.2),
            width=Inches(9), height=Inches(4)
        )
        desc_top = Inches(5.4)
    except:
        desc_top = Inches(1.2)
    
    # Add description
    desc_box = slide.shapes.add_textbox(
        Inches(0.5), desc_top, Inches(9), Inches(1.3)
    )
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = description
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(66, 66, 66)
    p.alignment = PP_ALIGN.LEFT


def create_presentation():
    """Create the PowerPoint presentation."""
    print("Creating Fabric Monitor demo presentation...")
    print(f"Output: {OUTPUT_PPTX}")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(SLIDES, 1):
        image_filename, title, description = slide_data
        print(f"[{idx}/{len(SLIDES)}] Creating slide: {title}")
        
        if image_filename is None:
            # Title slide
            add_title_slide(prs, title, description)
        else:
            # Content slide with image
            image_path = IMAGE_DIR / image_filename
            if not image_path.exists():
                print(f"  ⚠️  Warning: Image not found: {image_path}")
                # Still add slide without image
                add_content_slide(prs, None, title, description)
            else:
                add_content_slide(prs, image_path, title, description)
    
    # Save presentation
    print("\nSaving presentation...")
    prs.save(str(OUTPUT_PPTX))
    print(f"\n✓ Presentation created successfully: {OUTPUT_PPTX}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    if not IMAGE_DIR.exists():
        print(f"Error: Image directory not found: {IMAGE_DIR}")
        exit(1)
    
    try:
        create_presentation()
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        exit(1)
